'''
write_powerpoint.py

  pip3 install install python-pptx
'''


# Add pictures

from pptx import Presentation
from pptx.util import Inches, Pt

img_path = '/Users/carlos.aguilar/Documents/EF_Prism/F30d/charts/first_class_day_in_plot.png'

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)


left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = 'first_class_time_spent is H0=(17.08[0.00],21.25) and H1=(31.48[45.00],19.43)'
# tf.font.bold = True


left = top = Inches(1.5)
pic = slide.shapes.add_picture(img_path, left, top)

# add another slide
slide = prs.slides.add_slide(blank_slide_layout)


left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = 'first_class_time_spent is H0=(17.08[0.00],21.25) and H1=(31.48[45.00],19.43)'
# tf.font.bold = True


left = top = Inches(1.5)
pic = slide.shapes.add_picture(img_path, left, top)

prs.save('test.pptx')


