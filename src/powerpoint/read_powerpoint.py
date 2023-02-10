'''
	read_powerpoint.py

	https://python-pptx.readthedocs.io/en/latest/user/quickstart.html

	source ~/.bash_profile && python3 
	source ~/.bash_profile && python3 -m pip install python-pptx



	References:
	https://python-pptx.readthedocs.io/en/latest/api/presentation.html#presentation-function
	
	
'''

from pptx import Presentation

path_to_pptx_file = '/Users/carlos.aguilar/Google Drive/order/Machine Learning Part/ThesisPresentationVIVA/presentation backups/20.06.2021-PresentationDraft.pptx'
pptx = Presentation(path_to_pptx_file)

core_properties = pptx.core_properties

core_properties.author
core_properties.blob
core_properties.language
core_properties.content_type

core_properties.comments
core_properties.keywords
core_properties.title

core_properties.

# prs = Presentation()
# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]

#title.text = "Hello, World!"
#subtitle.text = "python-pptx was here!"

#prs.save('test.pptx')


sp = pptx.slides[0].shapes[0]



##########
from pptx import Presentation

path_to_pptx_file = '/Users/carlos.aguilar/Documents/EF_EVC_videos_hyperclass/lesson_plan/LP 9.1.2 - Ordering a meal.pptx'
path_to_pptx_file = '/Users/carlos.aguilar/Documents/EF_EVC_videos_hyperclass/lesson_plan/LP 9.2.3 - Negotiating.pptx'
pptx = Presentation(path_to_pptx_file)

core_properties = pptx.core_properties

core_properties.author
core_properties.blob
core_properties.language
core_properties.content_type

core_properties.comments
core_properties.keywords
core_properties.title



current_slide = pptx.slides[0]
shapes = current_slide.shapes

for shape in current_slide.shapes:
	if shape.has_text_frame:
		print(shape.text, '\n')
		input(">>>")

placeholders = current_slide.placeholders
placeholders.element
placeholders.part

current_slide = pptx.slides[1]
shapes = current_slide.shapes
for shape in current_slide.shapes:
	if shape.has_text_frame:
		print(shape.text, '\n')
		input(">>>")

current_slide = pptx.slides[2]
shapes = current_slide.shapes
for shape in current_slide.shapes:
	if shape.has_text_frame:
		print(shape.text, '\n')
		input(">>>")





# From https://www.geeksforgeeks.org/creating-and-updating-powerpoint-presentations-in-python-using-python-pptx/
# Example 3: PowerPoint (.pptx) file to Text (.txt) file conversion.
from pptx import Presentation
  
path_to_pptx_file = '/Users/carlos.aguilar/Documents/EF_EVC_videos_hyperclass/lesson_plan/LP 9.2.3 - Negotiating.pptx'
path_to_pptx_file = '/Users/carlos.aguilar/Documents/EF_EVC_videos_hyperclass/lesson_plan/LP 8.4.3 - Arranging a meeting.pptx'
path_to_pptx_file = '/Users/carlos.aguilar/Documents/EF_EVC_videos_hyperclass/lesson_plan/LP 9.1.2 - Ordering a meal.pptx'
pptx = Presentation(path_to_pptx_file)
  

# open file in write mode
# Write the text in the shapes
slides = []
text_filepath = path_to_pptx_file.replace('.pptx', '.txt')
with open(text_filepath, 'w') as fid:	
	for idx_slide, slide in enumerate(pptx.slides):
		shapes = []
		for idx_shape, shape in enumerate(slide.shapes, 0):
			if shape.has_text_frame:
				fid.write(f'slide[{idx_slide}][{idx_shape}]\n{shape.text}\n\n\n')
				shapes.append(shape.text)
		slides.append(shapes)
		print('\n\n')



slides[1][35]


# use paragraphs
text_filepath = path_to_pptx_file.replace('.pptx', '_paragraphs.txt')
# write text from powerpoint
# file into .txt file
with open(text_filepath, 'w') as fid:	
	for slide in pptx.slides:
		for shape in slide.shapes:
			if shape.has_text_frame:
				for paragraph in shape.text_frame.paragraphs:
					for run in paragraph.runs:
						fid.write(run.text+'\n')


# Add pictures

from pptx import Presentation
from pptx.util import Inches

img_path = 'monty-truth.png'

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(1)
pic = slide.shapes.add_picture(img_path, left, top)

left = Inches(5)
height = Inches(5.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

prs.save('test.pptx')


